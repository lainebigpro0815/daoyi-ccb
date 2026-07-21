import os, json, io, base64
from pathlib import Path
from fastapi import APIRouter, HTTPException, UploadFile, File
from fastapi.responses import HTMLResponse, FileResponse, JSONResponse

router = APIRouter(prefix="/api/projects/{project_id}/docs", tags=["docs"])

BASE_DIR = Path(__file__).parent.parent.parent.parent

OFFICE_EXTS = {".docx", ".xlsx", ".pptx", ".doc", ".xls", ".ppt"}
MARKDOWN_EXTS = {".md", ".markdown", ".txt"}
ALL_EXTS = OFFICE_EXTS | MARKDOWN_EXTS


def _project_docs_dir(project_id: int) -> Path:
    d = BASE_DIR / "docs" / "projects" / str(project_id)
    d.mkdir(parents=True, exist_ok=True)
    return d


@router.get("/tree")
def get_doc_tree(project_id: int):
    root = _project_docs_dir(project_id)

    def scan(dir_path: Path) -> list:
        items = []
        for entry in sorted(dir_path.iterdir(), key=lambda x: (not x.is_dir(), x.name)):
            if entry.name.startswith("."):
                continue
            rel = str(entry.relative_to(BASE_DIR)).replace("\\", "/")
            if entry.is_dir():
                children = scan(entry)
                items.append({"id": rel, "label": entry.name, "type": "folder", "children": children})
            elif entry.suffix in ALL_EXTS:
                ext = entry.suffix.lower()
                file_type = "word" if ext in (".docx", ".doc") else ("excel" if ext in (".xlsx", ".xls") else "ppt" if ext in (".pptx", ".ppt") else "markdown")
                items.append({"id": rel, "label": entry.name, "type": "file", "path": rel, "ext": ext, "file_type": file_type})
        return items

    return {"items": scan(root)}


@router.get("/read")
def read_doc(project_id: int, path: str):
    full_path = (BASE_DIR / path).resolve()
    if not str(full_path).startswith(str(BASE_DIR.resolve())):
        raise HTTPException(status_code=403)
    if not full_path.exists():
        raise HTTPException(status_code=404)

    ext = full_path.suffix.lower()
    content = ""
    html = ""

    if ext in MARKDOWN_EXTS:
        content = full_path.read_text(encoding="utf-8")

    elif ext == ".docx":
        try:
            import mammoth
            with open(full_path, "rb") as f:
                result = mammoth.convert_to_html(f)
                html = result.value
        except Exception as e:
            html = f"<p>预览不可用: {e}</p>"

    elif ext == ".xlsx":
        try:
            import openpyxl
            from openpyxl.utils import get_column_letter
            wb = openpyxl.load_workbook(full_path, data_only=True)
            sheets_html = ""
            for sname in wb.sheetnames:
                ws = wb[sname]
                sheets_html += f"<h4>{sname}</h4><table border='1' cellpadding='4' style='border-collapse:collapse;width:100%;margin-bottom:16px;font-size:13px;'>"
                for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row or 0, 100), values_only=True):
                    sheets_html += "<tr>" + "".join(f"<td>{str(v) if v is not None else ''}</td>" for v in row) + "</tr>"
                sheets_html += "</table>"
            html = sheets_html
        except Exception as e:
            html = f"<p>预览不可用: {e}</p>"

    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(full_path)
            slides_html = ""
            for i, slide in enumerate(prs.slides, 1):
                slides_html += f"<div style='margin-bottom:16px;padding:12px;background:#f5f7fa;border-radius:4px;'><strong>第{i}页</strong><br>"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slides_html += f"<p>{shape.text}</p>"
                slides_html += "</div>"
            html = slides_html
        except Exception as e:
            html = f"<p>预览不可用: {e}</p>"

    return {"path": path, "name": full_path.name, "content": content, "html": html, "ext": ext}


@router.get("/download")
def download_doc(project_id: int, path: str):
    full_path = (BASE_DIR / path).resolve()
    if not str(full_path).startswith(str(BASE_DIR.resolve())):
        raise HTTPException(status_code=403)
    if not full_path.exists():
        raise HTTPException(status_code=404)
    return FileResponse(path=str(full_path), filename=full_path.name)


@router.post("/upload")
async def upload_doc(project_id: int, file: UploadFile = File(...)):
    root = _project_docs_dir(project_id)
    dest = root / file.filename
    content = await file.read()
    dest.write_bytes(content)
    rel = str(dest.relative_to(BASE_DIR)).replace("\\", "/")
    return {"path": rel, "name": file.filename, "status": "ok"}


@router.post("/save")
def save_doc(project_id: int, data: dict):
    path = data.get("path", "")
    content = data.get("content", "")
    full_path = (BASE_DIR / path).resolve()
    if not str(full_path).startswith(str(BASE_DIR.resolve())):
        raise HTTPException(status_code=403)
    full_path.parent.mkdir(parents=True, exist_ok=True)
    full_path.write_text(content, encoding="utf-8")
    return {"status": "ok"}
